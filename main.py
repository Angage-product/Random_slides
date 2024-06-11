from pptx import Presentation
import os


def lire_pptx(fichier_pptx):
    # Ouvrir le fichier PowerPoint
    presentation = Presentation(fichier_pptx)

    return presentation

def count_slides(pres):
    # Parcourir les diapositives
    for i, diapo in enumerate(pres.slides):
        print(f"Diapositive {i + 1}")

def delete_slide(prs, slide):
    # Make dictionary with necessary information
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


def delete_slides_from_range(prs, from_, to):
    for i in range(to, from_, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]


def enregistrer_presentation(presentation, chemin_dossier, nom_fichier):
    # Crée le dossier s'il n'existe pas
    if not os.path.exists(chemin_dossier):
        os.makedirs(chemin_dossier)

    # Chemin complet pour le nouveau fichier
    chemin_fichier = os.path.join(chemin_dossier, nom_fichier)

    # Enregistre la présentation
    presentation.save(chemin_fichier)

def supprimer_diapositives(presentation, indices):
    # Supprime les diapositives en partant de la fin pour éviter les problèmes d'index
    for index in sorted(indices, reverse=True):
        xml_slides = presentation.slides._sldIdLst
        slide = xml_slides[index - 1]
        xml_slides.remove(slide)

# Nom du fichier PowerPoint à lire
fichier_pptx = 'pptx_in/demo_angage.pptx'
pres = lire_pptx(fichier_pptx)
count_slides(pres)
delete_slides_from_range(pres, 2, 4)
count_slides(pres)
enregistrer_presentation(pres, "pptx_out", "angage_demo_slides_deleted.pptx")